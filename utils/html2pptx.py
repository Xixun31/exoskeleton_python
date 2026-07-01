from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt
import re

# 1. 讀取你的 HTML 檔案
with open('encoder_imu.html', 'r', encoding='utf-8') as f:
    html_content = f.read()

# 2. 解析 HTML
soup = BeautifulSoup(html_content, 'html.parser')

# 3. 建立空白簡報
prs = Presentation()

# 4. 找出所有的投影片頁面
slides_divs = soup.find_all('div', class_='slide')

for i, slide_div in enumerate(slides_divs):
    # 第一頁：標題投影片
    if i == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        h1 = slide_div.find('h1')
        p = slide_div.find('p')
        
        if h1:
            slide.shapes.title.text = h1.text.strip()
        if p:
            slide.placeholders[1].text = p.text.strip()
            
    # 其他頁：內容投影片
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        h2 = slide_div.find('h2')
        
        if h2:
            slide.shapes.title.text = h2.text.strip()
            
        # 取得內文框架並清空預設文字
        tf = slide.placeholders[1].text_frame
        tf.clear()
        
        content_div = slide_div.find('div', class_='content')
        if not content_div:
            continue
            
        # --- 核心邏輯：針對你的 HTML Class 結構進行客製化抓取 ---
        # 尋找內容區塊中的 p 段落與各類 div 方塊
        for elem in content_div.find_all(['p', 'div']):
            classes = elem.get('class', [])
            
            # 處理一般段落 <p>
            if elem.name == 'p':
                p = tf.add_paragraph()
                p.text = elem.text.strip()
                p.level = 0
                
            # 處理方塊圖 <div class="box">
            elif 'box' in classes:
                p = tf.add_paragraph()
                # 替換換行符號，讓方塊內的文字在 PPT 呈現同一行
                clean_text = re.sub(r'\s+', ' ', elem.text.strip())
                p.text = f"■ 方塊節點：{clean_text}"
                p.level = 1
                
            # 處理連線文字 <div class="sys-line">
            elif 'sys-line' in classes:
                p = tf.add_paragraph()
                p.text = f"  ➔ 訊號線：{elem.text.strip()}"
                p.level = 2
                
            # 處理 While 迴圈標籤 <div class="loop-label">
            elif 'loop-label' in classes:
                p = tf.add_paragraph()
                p.text = f"【{elem.text.strip()}】"
                p.level = 1
                p.font.bold = True
                
            # 處理數學公式 <div class="math-box">
            elif 'math-box' in classes:
                p = tf.add_paragraph()
                p.text = f"公式：{elem.text.strip()}"
                p.level = 0
                p.font.bold = True
                
            # 處理 16-bit 資料視覺化 <div class="bit-container">
            elif 'bit-container' in classes:
                bits = []
                for bit in elem.find_all('div', class_='bit'):
                    val = bit.find('span', class_='bit-val').text.strip()
                    bits.append(val)
                
                p = tf.add_paragraph()
                p.text = f"資料結構：[{' | '.join(bits)}]"
                p.level = 1
                
            # 處理第 6 頁的角度換算區塊 <div class="col">
            elif 'col' in classes and not elem.find('div', class_='box'):
                texts = [s for s in elem.stripped_strings]
                if texts:
                    p = tf.add_paragraph()
                    p.text = " - ".join(texts)
                    p.level = 1

# 5. 儲存成 PPTX 檔案
output_filename = 'AS5145_Encoder_Presentation.pptx'
prs.save(output_filename)
print(f"✅ 轉換成功！已生成檔案：{output_filename}")