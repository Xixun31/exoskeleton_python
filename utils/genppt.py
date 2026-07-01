from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # 建立簡報物件 (預設為 4:3，可後續在 PPT 中改為 16:9)
    prs = Presentation()
    
    # 定義常用的投影片版面
    title_slide_layout = prs.slide_layouts[0] # 標題投影片
    blank_slide_layout = prs.slide_layouts[5] # 只有標題的空白投影片

    # ================= Slide 1: 標題頁 =================
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    title.text = "核心感測器讀取流程解析"
    subtitle.text = "AS5145 SPI 編碼器 ＆ IMU UART (MTData2) 架構"

    # ================= Slide 2: AS5145 硬體架構 =================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    slide2.shapes.title.text = "硬體架構：AS5145 Encoder (SPI)"
    
    # MCU Box
    mcu_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(2), Inches(3))
    mcu_box.text = "MCU (Master)"
    
    # Arrows & Text
    slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(3), Inches(3), Inches(0.2))
    slide2.shapes.add_textbox(Inches(3.5), Inches(2.6), Inches(2), Inches(0.5)).text = "MOSI (D4) ->"
    
    slide2.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.2), Inches(3.6), Inches(3), Inches(0.2))
    slide2.shapes.add_textbox(Inches(3.5), Inches(3.8), Inches(2), Inches(0.5)).text = "<- MISO (D5)"
    
    slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(4.2), Inches(3), Inches(0.2))
    slide2.shapes.add_textbox(Inches(3.5), Inches(4.4), Inches(2), Inches(0.5)).text = "SCLK (D3) ->"

    slide2.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(4.8), Inches(3), Inches(0.2))
    slide2.shapes.add_textbox(Inches(3.5), Inches(5.0), Inches(2), Inches(0.5)).text = "CS (D9) ->"

    # AS5145 Box
    as_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.5), Inches(2.5), Inches(3))
    as_box.text = "AS5145\n(12-bit Encoder)"
    as_box.fill.solid()
    as_box.fill.fore_color.rgb = RGBColor(211, 84, 0) # 橘色

    # ================= Slide 3: 主程式流程圖 =================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    slide3.shapes.title.text = "主程式迴圈 (Main Loop)"

    # Left Flow
    b1 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2.5), Inches(1))
    b1.text = "系統啟動 (Start)"
    slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(2), Inches(3.2), Inches(0.4), Inches(0.6))
    
    b2 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(2.5), Inches(1))
    b2.text = "初始化 IO 與 SPI\n(CS=1, 1MHz, 16-bit)"

    slide3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.8), Inches(3.2), Inches(0.8), Inches(0.4))

    # Right Flow (Loop)
    slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1.5), Inches(4), Inches(5))
    slide3.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(3), Inches(0.5)).text = "While (true) 迴圈"
    
    b3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.2), Inches(3), Inches(0.8))
    b3.text = "延遲 50ms (20Hz)"
    slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.8), Inches(3.1), Inches(0.3), Inches(0.5))

    b4 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.7), Inches(3), Inches(0.8))
    b4.text = "讀取編碼器與運算"
    b4.fill.solid()
    b4.fill.fore_color.rgb = RGBColor(211, 84, 0)

    slide3.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.8), Inches(4.6), Inches(0.3), Inches(0.5))

    b5 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5.2), Inches(3), Inches(0.8))
    b5.text = "印出數據 (printf)"

    # ================= Slide 4: SPI 讀取時序 =================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    slide4.shapes.title.text = "SPI 通訊讀取步驟 (read_position)"

    boxes = [
        ("CS 拉低 (0)\n啟動通訊", RGBColor(211, 84, 0)),
        ("Wait 10us\n給予反應時間", RGBColor(100, 100, 100)),
        ("寫入 0x00\n讀取 16-bit", RGBColor(0, 86, 179)),
        ("CS 拉高 (1)\n結束通訊", RGBColor(211, 84, 0))
    ]

    for i, (text, color) in enumerate(boxes):
        left = Inches(0.5 + i * 2.3)
        box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3), Inches(2), Inches(1.5))
        box.text = text
        box.fill.solid()
        box.fill.fore_color.rgb = color
        if i < 3:
            slide4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.05), Inches(3.6), Inches(0.2), Inches(0.2))

    # ================= Slide 5: 編碼器資料解析 =================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    slide5.shapes.title.text = "編碼器資料解析 (Data Parsing)"

    slide5.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(0.5)).text = "步驟 1: 原始 16-bit 資料 (raw_data)"
    slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(8), Inches(0.8)).text = "[15..12: 無效]  [11..0: 有效資料 D11~D0]"

    slide5.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(3.5), Inches(0.5), Inches(0.8))

    slide5.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5)).text = "步驟 2: 右移 3 位並 Mask 0x0FFF ( >> 3 & 0x0FFF )"
    b_step2 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5), Inches(8), Inches(0.8))
    b_step2.text = "[15..12: 補0]  [11..0: 對齊的有效資料 D11~D0]"
    b_step2.fill.solid()
    b_step2.fill.fore_color.rgb = RGBColor(0, 86, 179)

    # ================= Slide 6: 物理角度換算 =================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    slide6.shapes.title.text = "物理角度換算"

    tb = slide6.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tb.text = "將 12-bit (0 ~ 4095) 的乾淨數據，轉換為 0 ~ 360 度的實際物理角度。"
    
    formula = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(8), Inches(1.2))
    formula.text = "degree = (valid_data / 4096.0f) * 360.0f;"
    formula.fill.solid()
    formula.fill.fore_color.rgb = RGBColor(240, 240, 240)
    for p in formula.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.font.size = Pt(32)
        p.font.name = 'Courier New'

    # ================= Slide 7: IMU 架構 =================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    slide7.shapes.title.text = "硬體架構：IMU (UART 中斷模式)"

    mcu_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2.5), Inches(2), Inches(3))
    mcu_box.text = "MCU (Master)"
    
    slide7.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.2), Inches(3), Inches(3), Inches(0.2))
    slide7.shapes.add_textbox(Inches(3.5), Inches(2.6), Inches(2), Inches(0.5)).text = "<- RX (PA_10) 中斷"
    
    slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(3.2), Inches(4), Inches(3), Inches(0.2))
    slide7.shapes.add_textbox(Inches(3.5), Inches(4.2), Inches(2), Inches(0.5)).text = "TX (PA_9) ->"

    imu_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(2.5), Inches(2.5), Inches(3))
    imu_box.text = "IMU Sensor\n(MTData2 協定)\nBaud: 115200"
    imu_box.fill.solid()
    imu_box.fill.fore_color.rgb = RGBColor(46, 125, 50) # 綠色

    # ================= Slide 8: IMU 三步驟 =================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    slide8.shapes.title.text = "IMU 資料接收三步驟 (Rx ISR)"

    slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5)).text = "透過中斷 (Interrupt) 在背景自動收集資料，不阻塞主程式"

    imu_steps = [
        ("1. 找開頭\n(0xFA, 0xFF, 0x36)", RGBColor(0, 86, 179)),
        ("2. 看長度\n(預先確認 Payload 大小)", RGBColor(211, 84, 0)),
        ("3. 收資料\n(存入 Buffer 等待解析)", RGBColor(46, 125, 50))
    ]

    for i, (text, color) in enumerate(imu_steps):
        left = Inches(0.5 + i * 3)
        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(2.5), Inches(1.5))
        box.text = text
        box.fill.solid()
        box.fill.fore_color.rgb = color
        if i < 2:
            slide8.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(2.6), Inches(3.1), Inches(0.3), Inches(0.3))

    code_text = "switch (rx_state) {\n    case 0: if (byte == 0xFA) rx_state = 1; break;\n    case 3: expected_length = byte; rx_state = 4; break;\n    case 4: payload_buffer[idx++] = byte; break;\n}"
    code_box = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.5), Inches(8.5), Inches(2))
    code_box.text = code_text
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
    for p in code_box.text_frame.paragraphs:
        p.font.color.rgb = RGBColor(200, 200, 200)
        p.font.name = 'Consolas'
        p.alignment = PP_ALIGN.LEFT

    # ================= Slide 9: IMU 感測內容 =================
    slide9 = prs.slides.add_slide(blank_slide_layout)
    slide9.shapes.title.text = "IMU 感測器回傳內容解析"

    types = [
        ("姿態角 (Euler)", "Roll, Pitch, Yaw\n(0x2030)", RGBColor(0, 86, 179)),
        ("加速度 (Accel)", "aX, aY, aZ\n(0x4020)", RGBColor(211, 84, 0)),
        ("角速度 (Gyro)", "gX, gY, gZ\n(0x8020)", RGBColor(46, 125, 50))
    ]

    for i, (title_text, sub_text, color) in enumerate(types):
        left = Inches(0.5 + i * 3)
        box = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), Inches(2.5), Inches(2))
        box.text = f"{title_text}\n\n{sub_text}"
        box.fill.solid()
        box.fill.fore_color.rgb = color

    # ================= Slide 10: 防呆與轉換 =================
    slide10 = prs.slides.add_slide(blank_slide_layout)
    slide10.shapes.title.text = "資料轉換與防呆機制"

    b_swap = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(4.2), Inches(4))
    b_swap.text = "1. 資料重組 (Float Swap)\n\n將 4 個 Byte 組合回 32-bit 整數，再藉由記憶體複製 (memcpy) 還原成浮點數。\n\nuint32_t temp = (d[0]<<24) | ...\nfloat result;\nmemcpy(&result, &temp, 4);"
    b_swap.fill.solid()
    b_swap.fill.fore_color.rgb = RGBColor(0, 86, 179)

    b_safe = slide10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(2), Inches(4.2), Inches(4))
    b_safe.text = "2. 防護機制 (防止當機)\n\n若 UART 傳輸發生錯位，宣告的長度超出現有資料時，立刻中斷解析防越界。\n\nif (i + 3 + size > length) {\n    break; // 逃生出口\n}"
    b_safe.fill.solid()
    b_safe.fill.fore_color.rgb = RGBColor(211, 84, 0)

    # 儲存檔案
    prs.save('AS5145_IMU_Presentation.pptx')
    print("成功產生簡報：AS5145_IMU_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()